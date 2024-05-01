import os
from pptx import Presentation
from pptx.util import Pt
from pathlib import Path


def calculate_col_char_widths(table):
    # https://python-pptx.readthedocs.io/en/latest/dev/analysis/tbl-merge.html?highlight=merge%20cell#terminology

    col_widths = [[] for _ in range(len(table.columns))]

    cell_iterator = table.iter_cells()

    for _ in range(len(table.rows)):
        for col_i, _ in enumerate(table.columns):
            cell = next(cell_iterator)
            if cell.is_merge_origin or cell.is_spanned:
                col_width = 0
            else:
                col_width = len(cell.text)

            col_widths[col_i].append(col_width)

    return [max(widths) for widths in col_widths]


def table_to_markdown(table):
    table_lines = []
    col_widths = calculate_col_char_widths(table)

    cell_iterator = table.iter_cells()
    for row_i in range(len(table.rows)):
        if row_i == 1:
            # insert header separator
            row_text = " | ".join(
                [col_widths[col_i] * "-" for col_i in range(len(table.columns))]
            )
            table_lines.append(f"| {row_text} |")

        row_lines = []

        for col_i in range(len(table.columns)):
            cell = next(cell_iterator)
            col_width = col_widths[col_i]

            if cell.is_spanned or cell.is_merge_origin:
                cell_text = cell.text
            else:
                cell_text = cell.text.ljust(col_width)

            row_lines.append(cell_text)

        row_text = " | ".join(row_lines)
        table_lines.append(f"| {row_text} |")

    return table_lines


def to_markdown(input_file):
    presentation = Presentation(input_file)

    lines = []

    for presentation_i, slide in enumerate(presentation.slides):
        lines.append(f"Slide number: {presentation_i+1}")
        lines.append(f"Slide id: {slide.slide_id}")

        if slide.has_notes_slide:
            lines.append(slide.notes_slide())

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    lines.append(paragraph.text)

            # | Syntax      | Description |
            # | ----------- | ----------- |
            # | Header      | Title       |
            # | Paragraph   | Text        |
            if shape.has_table:
                lines += table_to_markdown(shape.table)

        for placeholder in slide.placeholders:
            lines.append(placeholder.text)
            # print(placeholder.shape_type)
            # print(placeholder.text)
            # print(placeholder.text_frame)
            # print(placeholder.has_chart)

        lines += "\n"

    return "\n".join(lines)


def _main():
    data_dir = Path(__file__).parents[2] / "data"
    output = to_markdown(input_file=data_dir / "sample.pptx")

    output_dir = Path(__file__).parents[2] / "output"
    output_dir.write_text(output)


if __name__ == "__main__":
    _main()
